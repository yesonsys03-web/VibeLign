# === ANCHOR: PPTX_RENDERER_START ===
from __future__ import annotations

import io

from vibelign.core.reporting_cli.docx_renderer import ReportRendererUnavailable
from vibelign.core.reporting_cli.font_sizes import ReportFontSizes
from vibelign.core.reporting_cli.fonts import ReportFonts, font_def
from vibelign.core.reporting_cli.models import ReportModel
from vibelign.core.reporting_cli.templates import meta_line
from vibelign.core.reporting_cli.themes import get_theme

try:
    import pptx  # noqa: F401
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# === ANCHOR: PPTX_RENDERER__SECTION_TEXT_START ===
def _section_text(section) -> str:
    lines: list[str] = []
    for block in section.blocks:
        if block.kind == "bullets":
            lines.extend(block.items)
        elif block.text:
            lines.append(block.text)
    return "\n".join(lines)
# === ANCHOR: PPTX_RENDERER__SECTION_TEXT_END ===


# === ANCHOR: PPTX_RENDERER_RENDER_PPTX_START ===
def render_pptx(
    model: ReportModel,
    theme: str = "classic",
    font_sizes: ReportFontSizes | None = None,
    fonts: ReportFonts | None = None,
) -> bytes:
    if not PPTX_AVAILABLE:
        raise ReportRendererUnavailable(
            "PPT 내보내기에 python-pptx 가 필요합니다. (pip install python-pptx)"
        )
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    accent = RGBColor.from_string(get_theme(theme).accent.lstrip("#"))
    heading_name = font_def(fonts.heading).office_name if fonts and fonts.heading else None
    body_name = font_def(fonts.body).office_name if fonts and fonts.body else None

    def _accent_title(shape) -> None:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = accent

    def _size_text(shape, value: int | None) -> None:
        if value is None or not shape.has_text_frame:
            return
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(value)

    def _font_text(shape, name: str | None) -> None:
        if name is None or not shape.has_text_frame:
            return
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = name
                rpr = run._r.get_or_add_rPr()
                for tag in ("a:ea", "a:cs"):
                    el = rpr.find(qn(tag))
                    if el is None:
                        el = rpr.makeelement(qn(tag), {})
                        rpr.append(el)
                    el.set("typeface", name)

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = model.title
    _accent_title(title_slide.shapes.title)
    _size_text(title_slide.shapes.title, font_sizes.title if font_sizes is not None else None)
    _font_text(title_slide.shapes.title, heading_name)
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = meta_line(model)
        _size_text(
            title_slide.placeholders[1],
            (font_sizes.meta or font_sizes.body) if font_sizes is not None else None,
        )
        _font_text(title_slide.placeholders[1], body_name)
    for section in model.sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.heading
        _accent_title(slide.shapes.title)
        _size_text(slide.shapes.title, font_sizes.heading if font_sizes is not None else None)
        _font_text(slide.shapes.title, heading_name)
        body = _section_text(section)
        if len(slide.placeholders) > 1 and body:
            slide.placeholders[1].text = body
            _size_text(slide.placeholders[1], font_sizes.body if font_sizes is not None else None)
            _font_text(slide.placeholders[1], body_name)
    buf = io.BytesIO()
# === ANCHOR: PPTX_RENDERER_RENDER_PPTX_END ===
    prs.save(buf)
    return buf.getvalue()
# === ANCHOR: PPTX_RENDERER_END ===
