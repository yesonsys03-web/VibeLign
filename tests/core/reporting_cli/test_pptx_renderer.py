import io
import pytest
from vibelign.core.reporting_cli.font_sizes import ReportFontSizes
from vibelign.core.reporting_cli.pptx_renderer import render_pptx, PPTX_AVAILABLE
from vibelign.core.reporting_cli.models import Block, ReportModel, Section


def _model():
    return ReportModel(
        title="예약 앱", report_type="work", date="2026-06-16",
        sections=[
            Section("개요", [Block(kind="summary", text="미용실 예약 앱")]),
            Section("핵심 내용", [Block(kind="bullets", items=["예약 캘린더", "알림 문자"])]),
        ],
    )


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx 미설치")
def test_render_pptx_opens_with_slide_per_section():
    data = render_pptx(_model())
    assert data[:2] == b"PK"
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    # 제목 슬라이드 1 + 섹션 2 = 3
    assert len(prs.slides) == 3
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    blob = "\n".join(all_text)
    assert "예약 앱" in blob
    assert "예약 캘린더" in blob


def test_render_pptx_missing_lib_raises_clear(monkeypatch):
    import vibelign.core.reporting_cli.pptx_renderer as mod
    monkeypatch.setattr(mod, "PPTX_AVAILABLE", False)
    with pytest.raises(mod.ReportRendererUnavailable):
        render_pptx(_model())


def test_pptx_theme_renders_valid_file():
    from vibelign.core.reporting_cli.models import Block, ReportModel, Section
    from vibelign.core.reporting_cli.pptx_renderer import render_pptx as _render

    m = ReportModel(title="t", report_type="work", date="d",
                    sections=[Section("개요", [Block(kind="bullets", items=["a"])])])
    data = _render(m, theme="pastel")
    assert data[:2] == b"PK"  # 유효한 pptx(zip)


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx 미설치")
def test_pptx_font_size_overrides_apply_to_title_heading_and_body():
    from pptx import Presentation

    data = render_pptx(_model(), font_sizes=ReportFontSizes(title=34, heading=22, body=16))
    prs = Presentation(io.BytesIO(data))
    title_shape = prs.slides[0].shapes.title
    section_title_shape = prs.slides[1].shapes.title
    body_shape = prs.slides[1].placeholders[1]
    assert title_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 34
    assert section_title_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 22
    assert body_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 16
